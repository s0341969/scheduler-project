from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "課堂打卡系統_專案簡報.pptx"
IMAGES = ROOT / "docs" / "images"

BRAND = RGBColor(16, 92, 99)
BRAND_DARK = RGBColor(12, 54, 60)
BRAND_LIGHT = RGBColor(219, 238, 238)
ACCENT = RGBColor(201, 136, 49)
ACCENT_LIGHT = RGBColor(250, 239, 221)
INK = RGBColor(29, 39, 51)
MUTED = RGBColor(97, 108, 120)
SURFACE = RGBColor(247, 244, 239)
SURFACE_ALT = RGBColor(240, 234, 226)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(220, 212, 200)
SUCCESS = RGBColor(40, 120, 76)
WARNING = RGBColor(176, 88, 59)


def set_background(slide, color: RGBColor = SURFACE) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, label: str) -> None:
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(25.4), Cm(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BRAND_DARK
    bar.line.fill.background()

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.3), Cm(0.78), Cm(4.6), Cm(0.8))
    tag.fill.solid()
    tag.fill.fore_color.rgb = BRAND
    tag.line.fill.background()
    tf = tag.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE


def add_title(slide, title: str, subtitle: str) -> None:
    add_top_bar(slide, "Classroom Attendance System")

    box = slide.shapes.add_textbox(Cm(1.4), Cm(1.9), Cm(18.6), Cm(1.6))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(25)
    run.font.bold = True
    run.font.color.rgb = INK

    sub = slide.shapes.add_textbox(Cm(1.45), Cm(3.0), Cm(20.0), Cm(0.9))
    p = sub.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(11.5)
    run.font.color.rgb = MUTED

    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(1.4), Cm(3.75), Cm(3.8), Cm(0.16))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Cm(1.3), Cm(18.1), Cm(22.9), Cm(0.5))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(8.8)
    run.font.color.rgb = MUTED


def add_body_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: float = 16,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bullets(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: float = 15,
    color: RGBColor = INK,
) -> None:
    box = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.space_after = Pt(7)
        p.font.name = "Microsoft JhengHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body: list[str] | str,
    fill_color: RGBColor = WHITE,
    title_color: RGBColor = BRAND_DARK,
) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = LINE

    add_body_text(slide, title, left + 0.35, top + 0.3, width - 0.7, 0.7, 15, title_color, True)

    if isinstance(body, list):
        add_bullets(slide, body, left + 0.3, top + 1.0, width - 0.6, height - 1.2, 12.8)
    else:
        add_body_text(slide, body, left + 0.3, top + 1.0, width - 0.6, height - 1.2, 12.5, MUTED)


def add_image(slide, image_name: str, left: float, top: float, width: float, height: float) -> None:
    path = IMAGES / image_name
    if path.exists():
        slide.shapes.add_picture(str(path), Cm(left), Cm(top), width=Cm(width), height=Cm(height))


def add_stat_card(slide, left: float, top: float, width: float, title: str, value: str, tone: RGBColor) -> None:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE

    add_body_text(slide, title, left + 0.35, top + 0.35, width - 0.7, 0.55, 10.5, MUTED)
    add_body_text(slide, value, left + 0.35, top + 0.95, width - 0.7, 0.9, 20, tone, True)


def add_timeline(slide, steps: list[tuple[str, str]]) -> None:
    base_left = 1.8
    base_top = 8.8
    gap = 5.9
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(2.3), Cm(base_top + 0.85), Cm(18.8), Cm(0.14))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND
    line.line.fill.background()

    for index, (title, body) in enumerate(steps):
        left = base_left + gap * index
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Cm(left), Cm(base_top), Cm(0.9), Cm(0.9))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if index % 2 == 0 else BRAND
        circle.line.fill.background()
        add_body_text(slide, str(index + 1), left, base_top + 0.1, 0.9, 0.5, 12, WHITE, True, PP_ALIGN.CENTER)

        add_body_text(slide, title, left - 0.2, base_top + 1.25, 2.2, 0.6, 11.5, INK, True, PP_ALIGN.CENTER)
        add_body_text(slide, body, left - 0.65, base_top + 1.9, 3.2, 1.2, 9.6, MUTED, False, PP_ALIGN.CENTER)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide, SURFACE)
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(1.2), Cm(1.1), Cm(10.2), Cm(1.1))
    hero.fill.solid()
    hero.fill.fore_color.rgb = BRAND
    hero.line.fill.background()
    add_body_text(slide, "Classroom Attendance System", 1.5, 1.43, 9.5, 0.5, 16, WHITE, True)
    add_body_text(slide, "課堂打卡系統", 1.4, 3.2, 9.5, 1.2, 28, INK, True)
    add_body_text(slide, "專案簡報\n從構想到規劃、實作、操作與未來展望", 1.45, 5.0, 9.8, 2.0, 17, MUTED)
    add_stat_card(slide, 1.4, 9.2, 3.0, "開發基礎", ".NET 9 MVC", BRAND)
    add_stat_card(slide, 4.7, 9.2, 3.0, "打卡方式", "登入 + QR", ACCENT)
    add_stat_card(slide, 8.0, 9.2, 3.0, "儲存模式", "JSON 單機版", SUCCESS)
    add_image(slide, "home.png", 12.8, 2.1, 10.8, 7.8)
    add_footer(slide, "版本整理日期：2026-05-14")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "簡報大綱", "以專題報告方式整理專案全貌")
    add_card(slide, 1.5, 4.4, 5.4, 4.5, "01 專案背景", ["構想來源", "使用情境", "目標與痛點"], BRAND_LIGHT)
    add_card(slide, 7.2, 4.4, 5.4, 4.5, "02 規劃與架構", ["功能規劃", "技術選型", "資料流程"], WHITE)
    add_card(slide, 12.9, 4.4, 5.4, 4.5, "03 實作成果", ["首頁總覽", "學生打卡", "管理後台"], BRAND_LIGHT)
    add_card(slide, 18.6, 4.4, 5.4, 4.5, "04 未來展望", ["功能擴充", "校務整合", "分析與治理"], WHITE)
    add_body_text(slide, "這份簡報聚焦在「為什麼做、怎麼做、做到什麼、接下來怎麼走」。", 1.6, 11.1, 20.5, 0.9, 15, MUTED)
    add_footer(slide, "Presentation Outline")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "專案背景與動機", "從教室現場出席管理痛點出發")
    add_card(slide, 1.4, 4.2, 7.0, 6.2, "現場痛點", [
        "紙本簽到容易代簽，且課後整理耗時",
        "臨時調課或加課時，管理方式不一致",
        "缺少即時統計與課後可追溯紀錄",
        "若沒有後台，老師難以快速開關打卡與匯出資料",
    ], WHITE)
    add_card(slide, 9.0, 4.2, 7.0, 6.2, "專案目標", [
        "提供教室現場即可執行的打卡工具",
        "讓學生打卡流程更簡短、身份更可信",
        "提供管理端可維護、可匯出、可追蹤的機制",
        "保留後續接校務系統與資料庫的擴充空間",
    ], BRAND_LIGHT)
    add_card(slide, 16.6, 4.2, 7.2, 6.2, "設計原則", [
        "先做單機可落地版本",
        "降低部署與使用門檻",
        "維持程式結構清楚，方便後續擴充",
        "在成本可控下增加本人驗證強度",
    ], WHITE)
    add_footer(slide, "核心想法：不是做最複雜，而是做最能在教室用起來的版本")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "需求規劃與功能地圖", "將需求拆成學生端、管理端與安全機制")
    add_card(slide, 1.4, 4.3, 7.2, 6.5, "學生端", [
        "學生登入",
        "掃描 QR Code",
        "進入課堂打卡頁",
        "送出出席紀錄",
    ], BRAND_LIGHT)
    add_card(slide, 9.0, 4.3, 7.2, 6.5, "管理端", [
        "課程維護與班級設定",
        "課堂時段管理",
        "開放 / 關閉打卡",
        "Excel 匯出與可疑紀錄複核",
    ], WHITE)
    add_card(slide, 16.6, 4.3, 7.2, 6.5, "風險控制", [
        "動態 QR Token",
        "班級 SSID 對應網段驗證",
        "重複打卡阻擋",
        "裝置指紋與來源資訊留存",
    ], BRAND_LIGHT)
    add_footer(slide, "功能規劃重點：流程完整、權責清楚、風險可追蹤")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "專案生命週期", "從想法到可用系統的落地過程")
    add_body_text(slide, "專案演進並非一次到位，而是按風險與使用價值逐步擴充。", 1.5, 4.5, 15.0, 0.8, 14, MUTED)
    add_timeline(slide, [
        ("概念形成", "定義教室現場\n需要的最小可用產品"),
        ("功能規劃", "拆出首頁、打卡、\n後台與匯出需求"),
        ("核心實作", "完成 MVC、JSON、\n課堂與出席模型"),
        ("安全補強", "加入學生登入、QR、\n網段驗證與可疑標記"),
    ])
    add_card(slide, 17.7, 4.5, 6.0, 3.8, "目前階段", [
        "已完成可示範、可操作、可維護版本",
        "已具備後續接校務系統的延伸基礎",
    ], ACCENT_LIGHT)
    add_footer(slide, "開發策略：先有穩定骨架，再往真實場域需求靠近")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "系統架構與技術選型", "兼顧部署簡單、維護性與擴充性")
    add_card(slide, 1.4, 4.3, 5.2, 6.4, "前端呈現", [
        "ASP.NET Core MVC",
        "Razor Views",
        "繁中教務風格 UI",
    ], WHITE)
    add_card(slide, 7.0, 4.3, 5.2, 6.4, "業務邏輯", [
        "AttendanceQueryService",
        "Admin / Student 驗證",
        "QR 與打卡規則判斷",
    ], BRAND_LIGHT)
    add_card(slide, 12.6, 4.3, 5.2, 6.4, "資料層", [
        "attendance-data.json",
        "JSON 持久化",
        "SemaphoreSlim 並發保護",
    ], WHITE)
    add_card(slide, 18.2, 4.3, 5.6, 6.4, "延伸能力", [
        "ClosedXML 匯出 Excel",
        "QRCoder 產生 QR Code",
        "Cookie Authentication",
        "可延伸接資料庫與 SSO",
    ], BRAND_LIGHT)
    add_footer(slide, "技術選型採取穩定、熟悉、可快速交付的方向")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "實作重點", "把需求轉成可運作且可維護的程式")
    add_bullets(slide, [
        "建立課程、課堂、出席紀錄三類核心模型，形成完整資料主軸",
        "首頁總覽依一年級至三年級與 101-310 班級分層顯示課堂資訊",
        "後台提供課程 / 課堂新增、編輯、刪除與打卡開關切換",
        "學生端打卡改採登入後提交，避免使用者自行修改學號姓名",
        "加入 QR Token 時效驗證、來源網段檢查與裝置指紋比對",
        "匯出報表納入打卡時間、IP、網段狀態與可疑原因，支援課後複核",
    ], 1.6, 4.4, 12.0, 10.6, 15.2)
    add_image(slide, "admin.png", 14.4, 4.4, 9.1, 6.5)
    add_footer(slide, "實作層面重視功能完整度與後續調整彈性")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "系統操作流程", "管理者與學生各自的使用路徑")
    add_card(slide, 1.4, 4.3, 10.6, 7.0, "管理者流程", [
        "1. 登入管理後台",
        "2. 建立課程並指定班級",
        "3. 建立課堂時段與主題",
        "4. 開啟 QR 打卡板供學生掃碼",
        "5. 課後檢視可疑紀錄並匯出 Excel",
    ], WHITE)
    add_card(slide, 12.5, 4.3, 10.7, 7.0, "學生流程", [
        "1. 先使用學生帳號登入",
        "2. 連上班級 SSID 對應網段",
        "3. 掃描課堂 QR Code 進入打卡頁",
        "4. 系統依登入身分完成打卡",
        "5. 若不符網段或重複打卡，系統會阻擋或標記",
    ], BRAND_LIGHT)
    add_footer(slide, "流程設計目的：減少人工確認步驟，同時提升打卡可信度")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "畫面成果：課堂總覽", "首頁以年級 / 班級快速定位課堂")
    add_image(slide, "home.png", 1.4, 4.0, 14.2, 10.1)
    add_card(slide, 16.0, 4.2, 7.4, 5.6, "畫面重點", [
        "儀表板顯示開放課堂數、累積打卡數與更新時間",
        "依一年級、二年級、三年級分區顯示班級課堂",
        "每張卡片可直接進入 QR 打卡板或學生登入",
    ], WHITE)
    add_card(slide, 16.0, 10.2, 7.4, 3.7, "管理價值", [
        "老師與教務端不必逐一翻找課堂",
        "班級層次清楚，更符合學校日常使用習慣",
    ], BRAND_LIGHT)
    add_footer(slide, "總覽頁是整個系統最常被使用的入口")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "畫面成果：學生登入與打卡", "縮短打卡路徑，同時強化本人性")
    add_image(slide, "login.png", 1.6, 4.5, 7.2, 5.2)
    add_image(slide, "checkin.png", 9.4, 4.5, 7.2, 5.2)
    add_card(slide, 17.2, 4.5, 6.0, 5.2, "設計重點", [
        "學號與姓名由登入身分帶入",
        "使用者只需填寫必要欄位",
        "流程短，適合課堂快速通行",
    ], WHITE)
    add_card(slide, 1.6, 10.4, 21.6, 3.0, "目前示範帳號", [
        "S1123001 / Student123!、S1123002 / Student123!、S1123003 / Student123!",
    ], ACCENT_LIGHT)
    add_footer(slide, "學生端核心目標：少輸入、少誤填、少偽造")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "畫面成果：QR 與後台管理", "支援課前開放、課後追蹤")
    add_image(slide, "qr-board.png", 1.5, 4.4, 8.0, 6.6)
    add_image(slide, "admin.png", 10.3, 4.4, 13.0, 6.6)
    add_body_text(slide, "QR 打卡板：用於投影與快速掃碼", 1.6, 11.3, 7.8, 0.7, 12.5, MUTED, True, PP_ALIGN.CENTER)
    add_body_text(slide, "管理後台：課程、課堂、Excel、可疑紀錄", 10.5, 11.3, 12.6, 0.7, 12.5, MUTED, True, PP_ALIGN.CENTER)
    add_footer(slide, "管理端同時覆蓋課前、課中、課後三個場景")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "安全與風險控制設計", "在教室場域下提升打卡可信度")
    add_card(slide, 1.4, 4.3, 7.2, 7.0, "已完成機制", [
        "學生帳號登入",
        "動態 QR Code 時效驗證",
        "班級 SSID 對應網段驗證",
        "同堂同學號防重複打卡",
        "IP、User-Agent、裝置指紋記錄",
    ], WHITE)
    add_card(slide, 9.1, 4.3, 7.2, 7.0, "可疑判定", [
        "來源 IP 不在允許網段",
        "缺少正常瀏覽器識別資訊",
        "同裝置近期被不同學號使用",
        "可疑資料在後台集中檢視",
    ], BRAND_LIGHT)
    add_card(slide, 16.8, 4.3, 7.0, 7.0, "設計評估", [
        "不依賴高成本硬體",
        "比純紙本或純表單更可信",
        "保留未來擴充生物辨識或卡務整合的空間",
    ], WHITE)
    add_footer(slide, "安全策略不是追求絕對，而是追求場域內最合理的成本效益")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "目前成果與價值", "專案已具備可展示、可操作、可延伸的基礎")
    add_stat_card(slide, 1.5, 4.4, 4.8, "課堂管理", "已可新增 / 編輯 / 關閉", BRAND)
    add_stat_card(slide, 6.7, 4.4, 4.8, "學生打卡", "登入 + QR + 驗證", ACCENT)
    add_stat_card(slide, 11.9, 4.4, 4.8, "資料留存", "JSON + Excel 匯出", SUCCESS)
    add_stat_card(slide, 17.1, 4.4, 6.0, "風險追蹤", "可疑紀錄複核", WARNING)
    add_card(slide, 1.5, 7.6, 10.8, 5.2, "對教務與老師的價值", [
        "減少點名、整理名單與匯出報表的人工作業",
        "打卡開關與 QR 發放方式更彈性",
        "出席資料可追溯，管理一致性更高",
    ], WHITE)
    add_card(slide, 12.8, 7.6, 10.3, 5.2, "對學生的價值", [
        "打卡流程簡短，不需反覆填寫個資",
        "課堂現場即掃即用，操作門檻低",
        "減少誤填、代填與重複提交情況",
    ], BRAND_LIGHT)
    add_footer(slide, "目前版本已經適合做專題展示與場域試行")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "未來展望", "從單機版走向校園級出席平台")
    add_card(slide, 1.4, 4.3, 7.0, 7.0, "短期優化", [
        "學生簽退與缺席統計",
        "管理者密碼變更 UI",
        "學生帳號與網段維護後台",
        "操作流程再簡化",
    ], WHITE)
    add_card(slide, 9.1, 4.3, 7.0, 7.0, "中期擴充", [
        "多管理者與角色權限",
        "校務單一登入整合",
        "多教室、多設備同步資料",
        "資料庫化與備份機制",
    ], BRAND_LIGHT)
    add_card(slide, 16.8, 4.3, 7.0, 7.0, "長期發展", [
        "出席分析儀表板",
        "請假 / 補點名 / 補簽流程",
        "校園卡或 NFC 整合",
        "與其他校務系統串接",
    ], WHITE)
    add_footer(slide, "未來方向：從打卡工具升級為完整出席管理平台")

    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "可增加功能提案", "依實際場域需求可以持續加值")
    add_card(slide, 1.4, 4.2, 5.2, 7.8, "教務管理類", [
        "班級與課程批次匯入",
        "學期課表排程匯入",
        "學生成績與出席關聯分析",
        "班級導師管理介面",
    ], WHITE)
    add_card(slide, 6.9, 4.2, 5.2, 7.8, "身份驗證類", [
        "校務帳號登入",
        "一次性驗證碼 OTP",
        "裝置綁定",
        "刷卡 / NFC / 學生證整合",
    ], BRAND_LIGHT)
    add_card(slide, 12.4, 4.2, 5.2, 7.8, "資料分析類", [
        "班級出席率圖表",
        "遲到 / 缺席統計",
        "可疑打卡熱點分析",
        "週報 / 月報自動產生",
    ], WHITE)
    add_card(slide, 17.9, 4.2, 5.2, 7.8, "整合服務類", [
        "Email / LINE 通知",
        "Google Calendar / 校務課表整合",
        "API 提供外部系統查詢",
        "雲端部署與多校區支援",
    ], BRAND_LIGHT)
    add_footer(slide, "功能增加方向可分為：治理、驗證、分析、整合四條路線")

    slide = prs.slides.add_slide(blank)
    set_background(slide, BRAND_DARK)
    add_body_text(slide, "結論", 1.5, 3.0, 4.0, 0.7, 15, RGBColor(220, 230, 229), True)
    add_body_text(slide, "課堂打卡系統已完成從構想到可運作雛型的落地。", 1.5, 4.5, 15.5, 1.2, 25, WHITE, True)
    add_body_text(slide, "目前版本重點在於：\n1. 教室現場可用\n2. 管理流程完整\n3. 本人驗證強度提升\n4. 後續擴充空間明確", 1.6, 7.1, 10.8, 3.2, 16, RGBColor(226, 234, 234))
    add_body_text(slide, "下一步可依學校需求，逐步往校務整合、資料治理與分析報表發展。", 1.6, 11.0, 11.4, 1.3, 15, RGBColor(226, 234, 234))
    add_image(slide, "home.png", 14.0, 4.2, 9.0, 6.5)
    add_footer(slide, "課堂打卡系統｜Project Presentation")

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
